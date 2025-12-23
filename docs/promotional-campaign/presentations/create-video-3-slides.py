#!/usr/bin/env python3
"""
Generate PowerPoint slides for Video 3: One Server, Unlimited Scenarios

This script creates a presentation showing how Scenarist makes "hard" and "impossible"
testing scenarios trivially easy by mocking server-side HTTP calls.

Usage:
    python3 create-video-3-slides.py

Output:
    video-03-one-server-unlimited-scenarios.pptx

Requirements:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette (matching Video 1 & 2)
DARK_BG = RGBColor(24, 24, 27)  # zinc-900
WHITE = RGBColor(255, 255, 255)
ZINC_400 = RGBColor(161, 161, 170)
YELLOW = RGBColor(245, 158, 11)  # amber-500
GREEN = RGBColor(34, 197, 94)  # green-500
RED = RGBColor(239, 68, 68)  # red-500
BLUE = RGBColor(59, 130, 246)  # blue-500

# Service colors
USER_BLUE = RGBColor(59, 130, 246)
INVENTORY_GREEN = RGBColor(34, 197, 94)
SHIPPING_PURPLE = RGBColor(168, 85, 247)  # purple-500


def add_dark_slide(prs):
    """Add a slide with dark background."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    return slide


def add_title(slide, text, top=Inches(0.5), font_size=54, color=WHITE):
    """Add a title text box."""
    left = Inches(0.5)
    width = Inches(12.33)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_subtitle(slide, text, top=Inches(1.5), font_size=28, color=ZINC_400):
    """Add a subtitle text box."""
    left = Inches(0.5)
    width = Inches(12.33)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER


def add_bullet_point(slide, text, left, top, color=WHITE, bullet=""):
    """Add a bullet point."""
    width = Inches(10)
    height = Inches(0.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{bullet} {text}" if bullet else text
    p.font.size = Pt(24)
    p.font.color.rgb = color


def add_code_block(slide, code, left, top, width, height, font_size=14):
    """Add a code block with monospace font."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False

    # Add background
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(39, 39, 42)  # zinc-800

    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(167, 243, 208)  # green-200
    p.font.name = "Menlo"


def create_presentation():
    """Create the Video 3 presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============================================================================
    # SLIDE 1: Title
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "One Server, Unlimited Scenarios", top=Inches(2.5), font_size=54)
    add_subtitle(slide, "Video 3: How Scenarist Makes Testing Easy", top=Inches(3.5), font_size=32, color=ZINC_400)

    # ============================================================================
    # SLIDE 2: Recap - Testing Problem Table
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Remember the Testing Problem Table?", top=Inches(0.5), font_size=44)

    # Problem summary
    add_bullet_point(slide, "Happy path", Inches(1.5), Inches(1.8), GREEN, "Easy:")
    add_bullet_point(slide, "User tiers (pro/free)", Inches(1.5), Inches(2.4), YELLOW, "Annoying:")
    add_bullet_point(slide, "Offer states, service errors", Inches(1.5), Inches(3.0), YELLOW, "Hard:")
    add_bullet_point(slide, "Offer ends during checkout", Inches(1.5), Inches(3.6), RED, "Impossible:")
    add_bullet_point(slide, "50 parallel tests", Inches(1.5), Inches(4.2), RED, "Impossible:")

    add_subtitle(slide, "Let's fix that.", top=Inches(5.5), font_size=36, color=YELLOW)

    # ============================================================================
    # SLIDE 3: PayFlow Architecture - Three Services
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "PayFlow's Architecture", top=Inches(0.4), font_size=44)
    add_subtitle(slide, "Three backend services - all server-side HTTP calls", top=Inches(1.0), font_size=24, color=ZINC_400)

    # Browser box
    browser_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(2), Inches(1))
    browser_box.fill.solid()
    browser_box.fill.fore_color.rgb = ZINC_400
    browser_box.line.fill.background()
    tf = browser_box.text_frame
    tf.paragraphs[0].text = "Browser"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Next.js Server box
    nextjs_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.5), Inches(2.5), Inches(1))
    nextjs_box.fill.solid()
    nextjs_box.fill.fore_color.rgb = WHITE
    nextjs_box.line.fill.background()
    tf = nextjs_box.text_frame
    tf.paragraphs[0].text = "Next.js Server"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Service boxes
    services = [
        ("User Service", "/users/current", USER_BLUE, Inches(4.2)),
        ("Inventory Service", "/inventory", INVENTORY_GREEN, Inches(5.0)),
        ("Shipping Service", "/shipping", SHIPPING_PURPLE, Inches(5.8)),
    ]

    for name, endpoint, color, top in services:
        svc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), top, Inches(3.5), Inches(0.7))
        svc_box.fill.solid()
        svc_box.fill.fore_color.rgb = color
        svc_box.line.fill.background()
        tf = svc_box.text_frame
        tf.paragraphs[0].text = f"{name}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Key point at bottom
    add_subtitle(slide, "Your browser never talks to these services directly.", top=Inches(6.5), font_size=24, color=YELLOW)

    # ============================================================================
    # SLIDE 4: The Core Insight
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "The Core Insight", top=Inches(2), font_size=54)
    add_subtitle(slide, "If we can intercept those server-side calls,", top=Inches(3.2), font_size=32, color=WHITE)
    add_subtitle(slide, "we can return whatever we want.", top=Inches(4), font_size=32, color=YELLOW)

    # ============================================================================
    # SLIDE 5: What This Means
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Control Every Response", top=Inches(0.8), font_size=44)

    add_bullet_point(slide, "Pro user? Return { tier: 'pro' }", Inches(2), Inches(2.2), GREEN)
    add_bullet_point(slide, "Free user? Return { tier: 'free' }", Inches(2), Inches(3.0), GREEN)
    add_bullet_point(slide, "Offer ended? Return { quantity: 0 }", Inches(2), Inches(3.8), GREEN)
    add_bullet_point(slide, "Shipping down? Return 500 error", Inches(2), Inches(4.6), GREEN)

    add_subtitle(slide, "No database changes. No test accounts. Just responses.", top=Inches(5.8), font_size=28, color=YELLOW)

    # ============================================================================
    # SLIDE 6: Scenarist Introduction
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Scenarist", top=Inches(1.5), font_size=64, color=YELLOW)
    add_subtitle(slide, "Intercepts server-side HTTP calls.", top=Inches(3), font_size=32, color=WHITE)
    add_subtitle(slide, "Returns whatever you specify.", top=Inches(4), font_size=32, color=GREEN)

    # ============================================================================
    # SLIDE 7: Architecture Overview
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Framework-Agnostic Architecture", top=Inches(0.4), font_size=40)

    # Diagram boxes
    box_width = Inches(5)
    box_height = Inches(1)
    center_left = Inches(4.166)

    # Test box
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_left, Inches(1.3), box_width, box_height)
    box1.fill.solid()
    box1.fill.fore_color.rgb = BLUE
    box1.line.fill.background()
    tf = box1.text_frame
    p = tf.paragraphs[0]
    p.text = "Your Test"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    add_subtitle(slide, "↓", top=Inches(2.2), font_size=36, color=WHITE)

    # Scenarist Core box
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_left, Inches(2.8), box_width, box_height)
    box2.fill.solid()
    box2.fill.fore_color.rgb = YELLOW
    box2.line.fill.background()
    tf = box2.text_frame
    p = tf.paragraphs[0]
    p.text = "Scenarist Core"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    add_subtitle(slide, "↓", top=Inches(3.7), font_size=36, color=WHITE)

    # Adapter box
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_left, Inches(4.3), box_width, box_height)
    box3.fill.solid()
    box3.fill.fore_color.rgb = GREEN
    box3.line.fill.background()
    tf = box3.text_frame
    p = tf.paragraphs[0]
    p.text = "Framework Adapter"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    add_subtitle(slide, "↓", top=Inches(5.2), font_size=36, color=WHITE)

    # MSW box
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_left, Inches(5.8), box_width, box_height)
    box4.fill.solid()
    box4.fill.fore_color.rgb = RGBColor(168, 85, 247)  # purple
    box4.line.fill.background()
    tf = box4.text_frame
    p = tf.paragraphs[0]
    p.text = "MSW (Interception)"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Annotations
    add_bullet_point(slide, "← Framework-agnostic", Inches(9.5), Inches(3), ZINC_400)
    add_bullet_point(slide, "← Express, Next.js, etc", Inches(9.5), Inches(4.5), ZINC_400)

    # ============================================================================
    # SLIDE 8: Key Message - Adapters
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "One Pattern, Any Framework", top=Inches(2.5), font_size=48)
    add_subtitle(slide, "Express today? Next.js tomorrow?", top=Inches(3.8), font_size=32, color=WHITE)
    add_subtitle(slide, "Only the adapter changes.", top=Inches(4.6), font_size=32, color=YELLOW)

    # ============================================================================
    # SLIDE 9: Live Demo Setup
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Live Demo", top=Inches(0.5), font_size=54, color=YELLOW)

    # Two terminals illustration
    add_subtitle(slide, "Two terminals:", top=Inches(1.8), font_size=28, color=WHITE)

    # Terminal boxes side by side
    term1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.8), Inches(4), Inches(2.5))
    term1.fill.solid()
    term1.fill.fore_color.rgb = RGBColor(39, 39, 42)
    term1.line.color.rgb = ZINC_400
    tf = term1.text_frame
    tf.paragraphs[0].text = "Terminal 1"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = ZINC_400
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "\n$ pnpm dev"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN
    p.font.name = "Menlo"
    p = tf.add_paragraph()
    p.text = "Next.js running..."
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Menlo"

    term2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.8), Inches(4), Inches(2.5))
    term2.fill.solid()
    term2.fill.fore_color.rgb = RGBColor(39, 39, 42)
    term2.line.color.rgb = ZINC_400
    tf = term2.text_frame
    tf.paragraphs[0].text = "Terminal 2"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = ZINC_400
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "\n$ pnpm inventory"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN
    p.font.name = "Menlo"
    p = tf.add_paragraph()
    p.text = "Backend services..."
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Menlo"

    add_subtitle(slide, "Watch Terminal 2 - it should stay silent.", top=Inches(6), font_size=24, color=YELLOW)

    # ============================================================================
    # SLIDE 10: Demo Test 1 - Pro User
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Test 1: Pro User Discount", top=Inches(0.5), font_size=40)

    code = '''test("pro user sees 20% discount", async ({ page, switchScenario }) => {
  await switchScenario("default");
  await page.goto("/products/1");
  await expect(page.getByText("20% off")).toBeVisible();
});'''
    add_code_block(slide, code, Inches(1), Inches(1.8), Inches(11), Inches(2.5), font_size=18)

    add_subtitle(slide, "User Service returns { tier: 'pro' }", top=Inches(5), font_size=28, color=GREEN)
    add_subtitle(slide, "Discount applied.", top=Inches(5.8), font_size=28, color=WHITE)

    # ============================================================================
    # SLIDE 11: Demo Test 2 - Free User
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Test 2: Free User Full Price", top=Inches(0.5), font_size=40)

    code = '''test("free user sees full price", async ({ page, switchScenario }) => {
  await switchScenario("freeUser");
  await page.goto("/products/1");
  await expect(page.getByText("$99.99")).toBeVisible();
});'''
    add_code_block(slide, code, Inches(1), Inches(1.8), Inches(11), Inches(2.5), font_size=18)

    add_subtitle(slide, "User Service returns { tier: 'free' }", top=Inches(5), font_size=28, color=GREEN)
    add_subtitle(slide, "No test account needed. Just a different response.", top=Inches(5.8), font_size=28, color=YELLOW)

    # ============================================================================
    # SLIDE 12: Demo Test 3 - Offer Ended
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Test 3: Offer Ended", top=Inches(0.5), font_size=40)

    code = '''test("offer ended shows sold out", async ({ page, switchScenario }) => {
  await switchScenario("offerEnded");
  await page.goto("/products/1");
  await expect(page.getByText("Sold Out")).toBeVisible();
});'''
    add_code_block(slide, code, Inches(1), Inches(1.8), Inches(11), Inches(2.5), font_size=18)

    add_subtitle(slide, "Inventory Service returns { quantity: 0 }", top=Inches(5), font_size=28, color=GREEN)
    add_subtitle(slide, "No database editing. Just a different response.", top=Inches(5.8), font_size=28, color=YELLOW)

    # ============================================================================
    # SLIDE 13: Demo Test 4 - Shipping Service Down
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Test 4: Shipping Service Down", top=Inches(0.5), font_size=40)

    code = '''test("handles shipping errors gracefully", async ({ page, switchScenario }) => {
  await switchScenario("shippingServiceDown");
  await page.goto("/checkout");
  await expect(page.getByText("Unable to load shipping")).toBeVisible();
});'''
    add_code_block(slide, code, Inches(1), Inches(1.8), Inches(11), Inches(2.5), font_size=18)

    add_subtitle(slide, "Shipping Service returns 500 error", top=Inches(5), font_size=28, color=RED)
    add_subtitle(slide, "Test your error handling without breaking anything.", top=Inches(5.8), font_size=28, color=YELLOW)

    # ============================================================================
    # SLIDE 14: Zero Requests
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Backend Services Terminal:", top=Inches(1.5), font_size=40)

    # Big zero
    zero_box = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(5), Inches(2))
    tf = zero_box.text_frame
    p = tf.paragraphs[0]
    p.text = "0 Requests"
    p.font.size = Pt(72)
    p.font.color.rgb = YELLOW
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_subtitle(slide, "Scenarist intercepted everything.", top=Inches(5.2), font_size=32, color=GREEN)
    add_subtitle(slide, "The real services are running, but we never hit them.", top=Inches(6), font_size=24, color=WHITE)

    # ============================================================================
    # SLIDE 15: Summary - What We Just Did
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Four Scenarios", top=Inches(0.8), font_size=48)

    add_bullet_point(slide, "Pro user discount (was: Annoying)", Inches(2), Inches(2.2), GREEN, "1.")
    add_bullet_point(slide, "Free user full price (was: Annoying)", Inches(2), Inches(3.0), GREEN, "2.")
    add_bullet_point(slide, "Offer ended (was: Hard)", Inches(2), Inches(3.8), GREEN, "3.")
    add_bullet_point(slide, "Shipping service down (was: Hard)", Inches(2), Inches(4.6), GREEN, "4.")

    add_subtitle(slide, "Now they're just... tests.", top=Inches(5.8), font_size=36, color=YELLOW)

    # ============================================================================
    # SLIDE 16: Parallel Test Isolation
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "Parallel Test Isolation", top=Inches(0.5), font_size=44)
    add_subtitle(slide, "What about running 50 tests at once?", top=Inches(1.3), font_size=24, color=ZINC_400)

    # Test ID diagram
    tests = [
        ("Test A", "abc123", "proUser", "20% discount", USER_BLUE),
        ("Test B", "def456", "freeUser", "Full price", INVENTORY_GREEN),
        ("Test C", "ghi789", "offerEnded", "Sold Out", SHIPPING_PURPLE),
        ("Test D", "jkl012", "shippingDown", "Error handling", RED),
    ]

    for i, (test, tid, scenario, result, color) in enumerate(tests):
        top = Inches(2.2 + i * 0.9)
        # Test box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), top, Inches(2.5), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.paragraphs[0].text = f"{test} ({tid[:6]})"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Arrow
        add_bullet_point(slide, "→", Inches(3.6), top, WHITE)

        # Scenario
        add_bullet_point(slide, scenario, Inches(4.2), top, ZINC_400)

        # Arrow
        add_bullet_point(slide, "→", Inches(6.8), top, WHITE)

        # Result
        add_bullet_point(slide, result, Inches(7.4), top, color)

    add_subtitle(slide, "Same server. Same endpoints. Different responses.", top=Inches(6.2), font_size=28, color=YELLOW)

    # ============================================================================
    # SLIDE 17: Key Point - Test ID Header
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "x-scenarist-test-id", top=Inches(2), font_size=48, color=YELLOW)
    add_subtitle(slide, "Every request includes this header.", top=Inches(3.2), font_size=28, color=WHITE)
    add_subtitle(slide, "Scenarist looks up the scenario for that specific test.", top=Inches(4), font_size=28, color=WHITE)
    add_subtitle(slide, "50 tests in parallel. Zero conflicts.", top=Inches(5.2), font_size=32, color=GREEN)

    # ============================================================================
    # SLIDE 18: Tease - The Killer Scenario
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "But What About...", top=Inches(0.8), font_size=44)
    add_subtitle(slide, '"Offer ends during checkout"', top=Inches(1.8), font_size=36, color=RED)
    add_subtitle(slide, "The one that was truly impossible?", top=Inches(2.6), font_size=24, color=ZINC_400)

    code = '''// Response Sequences
{
  url: "/inventory/1",
  sequence: [
    { quantity: 15 },  // First call: available
    { quantity: 0 },   // Second call: sold out
  ]
}'''
    add_code_block(slide, code, Inches(2), Inches(3.5), Inches(9), Inches(2.8), font_size=18)

    add_subtitle(slide, "That's the next video.", top=Inches(6.5), font_size=28, color=YELLOW)

    # ============================================================================
    # SLIDE 19: End Card
    # ============================================================================
    slide = add_dark_slide(prs)
    add_title(slide, "One Server, Unlimited Scenarios", top=Inches(2.5), font_size=48)
    add_subtitle(slide, "github.com/citypaul/scenarist", top=Inches(3.8), font_size=28, color=YELLOW)
    add_subtitle(slide, "Next: Response Sequences", top=Inches(5), font_size=24, color=ZINC_400)

    # Save the presentation
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(script_dir, "video-03-one-server-unlimited-scenarios.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    create_presentation()
