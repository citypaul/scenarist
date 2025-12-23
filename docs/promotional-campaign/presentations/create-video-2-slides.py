#!/usr/bin/env python3
"""
Generate PowerPoint slides for Video 2: Meet PayFlow - A Real Payment App

Run: python3 create-video-2-slides.py
Output: video-02-meet-payflow.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(24, 24, 27)  # zinc-900
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(239, 68, 68)  # red-500
YELLOW = RGBColor(245, 158, 11)  # amber-500
GREEN = RGBColor(34, 197, 94)  # green-500
BLUE = RGBColor(59, 130, 246)  # blue-500
GRAY = RGBColor(161, 161, 170)  # zinc-400
USER_SERVICE = RGBColor(99, 102, 241)  # indigo-500
INVENTORY_GREEN = RGBColor(16, 185, 129)  # emerald-500
SHIPPING_BLUE = RGBColor(14, 165, 233)  # sky-500


def add_dark_slide(prs):
    """Add a slide with dark background."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    # Add dark background shape
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide


def add_title(slide, text, top=Inches(2.5), font_size=60, color=WHITE):
    """Add centered title text."""
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_subtitle(slide, text, top=Inches(4), font_size=32, color=GRAY):
    """Add centered subtitle text."""
    txBox = slide.shapes.add_textbox(Inches(1), top, Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_code_block(slide, code, left, top, width, height):
    """Add a code block with monospace font."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(39, 39, 42)  # zinc-800
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Menlo"
    p.font.color.rgb = RGBColor(228, 228, 231)  # zinc-200
    return shape


def add_bullet_point(slide, text, left, top, color=WHITE, icon=""):
    """Add a bullet point with optional icon."""
    txBox = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon}  {text}" if icon else text
    p.font.size = Pt(28)
    p.font.color.rgb = color
    return txBox


def add_service_box(slide, name, description, color, left, top, width=Inches(3.5), height=Inches(2.5)):
    """Add a service box with name and description."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), width - Inches(0.4), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = f"\n{description}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(228, 228, 231)
    p.alignment = PP_ALIGN.CENTER
    return box


# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Meet PayFlow", top=Inches(2.8), font_size=72)
add_subtitle(slide, "A Real Payment App", top=Inches(4.2), font_size=40, color=GRAY)

# ============================================================================
# SLIDE 2: What We're Building
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Before we test...", top=Inches(2.5), font_size=48)
add_title(slide, "Let's see what we're testing.", top=Inches(3.8), font_size=48, color=YELLOW)

# ============================================================================
# SLIDE 3: Tech Stack
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Tech Stack", top=Inches(0.8), font_size=44)
add_bullet_point(slide, "Next.js 16 (App Router)", Inches(3), Inches(2), WHITE)
add_bullet_point(slide, "TypeScript", Inches(3), Inches(2.7), WHITE)
add_bullet_point(slide, "Tailwind CSS + shadcn/ui", Inches(3), Inches(3.4), WHITE)
add_bullet_point(slide, "Three backend services", Inches(3), Inches(4.1), YELLOW)
add_subtitle(slide, "Nothing exotic. The kind of app you'd actually build.", top=Inches(5.5), font_size=28, color=GRAY)

# ============================================================================
# SLIDE 4: The Three Services
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Three Backend Services", top=Inches(0.5), font_size=44)

# Service boxes
box_width = Inches(3.8)
box_height = Inches(3)
gap = Inches(0.4)
start_left = Inches(0.8)

# User Service
add_service_box(slide, "User Service", "User Tier\nPro/Free\n/users/current", USER_SERVICE,
                start_left, Inches(1.8), box_width, box_height)

# Inventory Service
add_service_box(slide, "Inventory Service", "Offer Availability\nQuantity/Reserved\n/inventory/:id", INVENTORY_GREEN,
                start_left + box_width + gap, Inches(1.8), box_width, box_height)

# Shipping Service
add_service_box(slide, "Shipping Service", "Delivery Options\nRates & Times\n/shipping", SHIPPING_BLUE,
                start_left + 2 * (box_width + gap), Inches(1.8), box_width, box_height)

add_subtitle(slide, "All server-side HTTP calls. Browser → Next.js → Services.", top=Inches(5.5), font_size=28, color=YELLOW)

# ============================================================================
# SLIDE 5: The Key Point - Server-Side
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Key Point", top=Inches(1), font_size=48)
add_subtitle(slide, "Your browser never talks to these services directly.", top=Inches(2.2), font_size=32, color=GRAY)
add_title(slide, "Next.js makes server-side HTTP calls.", top=Inches(3.5), font_size=36, color=WHITE)

# Architecture diagram
arch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.5), Inches(10.333), Inches(1.5))
arch_box.fill.solid()
arch_box.fill.fore_color.rgb = RGBColor(39, 39, 42)
arch_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.7), Inches(4.7), Inches(9.933), Inches(1.1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Browser → Next.js Server → User / Inventory / Shipping Services"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "This is 100% mockable.", top=Inches(6.3), font_size=24, color=GREEN)

# ============================================================================
# SLIDE 6: Code - User Service
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "User Service", top=Inches(0.5), font_size=36, color=USER_SERVICE)

code = """// Server-side call to User Service
const response = await fetch("http://localhost:3001/users/current");
const user = await response.json();
// Returns: { id, email, name, tier }

// tier: "free" | "basic" | "pro" | "enterprise"
// Pro users get 20% discount"""

add_code_block(slide, code, Inches(1.5), Inches(1.8), Inches(10), Inches(3))
add_subtitle(slide, "Server fetches user tier for pricing decisions.", top=Inches(5.5), font_size=28, color=WHITE)

# ============================================================================
# SLIDE 7: Code - Inventory Service
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Inventory Service", top=Inches(0.5), font_size=36, color=INVENTORY_GREEN)

code = """// Server-side call to Inventory Service
const response = await fetch(
  `http://localhost:3001/inventory/${productId}`
);
const data = await response.json();
// Returns: { id, productId, quantity, reserved }

// quantity: 0 = offer ended, 3 = limited spots"""

add_code_block(slide, code, Inches(1.5), Inches(1.8), Inches(10), Inches(3.2))
add_subtitle(slide, "Server fetches offer availability. No test mode.", top=Inches(5.5), font_size=28, color=WHITE)

# ============================================================================
# SLIDE 8: Code - Shipping Service
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Shipping Service", top=Inches(0.5), font_size=36, color=SHIPPING_BLUE)

code = """// Server-side call to Shipping Service
const response = await fetch("http://localhost:3001/shipping");
const options = await response.json();
// Returns: [
//   { id: "standard", price: 5.99, estimatedDays: "5-7" },
//   { id: "express", price: 14.99, estimatedDays: "2-3" },
//   { id: "overnight", price: 29.99, estimatedDays: "1" }
// ]"""

add_code_block(slide, code, Inches(1.5), Inches(1.8), Inches(10), Inches(3.5))
add_subtitle(slide, "Server fetches delivery options and rates.", top=Inches(5.8), font_size=28, color=WHITE)

# ============================================================================
# SLIDE 9: Live Demo - Setup
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Two Terminals", top=Inches(0.8), font_size=44)

# Terminal boxes
term_width = Inches(5.5)
term_height = Inches(2.8)
term_gap = Inches(0.5)
term_start = Inches(1.2)

# Terminal 1: Next.js
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, term_start, Inches(2), term_width, term_height)
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(39, 39, 42)
box1.line.fill.background()

txBox = slide.shapes.add_textbox(term_start + Inches(0.2), Inches(2.2), term_width - Inches(0.4), Inches(2.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Next.js"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p = tf.add_paragraph()
p.text = "pnpm dev"
p.font.size = Pt(18)
p.font.name = "Menlo"
p.font.color.rgb = GREEN
p = tf.add_paragraph()
p.text = "localhost:3000"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p = tf.add_paragraph()
p.text = "\nYour app"
p.font.size = Pt(16)
p.font.color.rgb = GRAY

# Terminal 2: Backend Services
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, term_start + term_width + term_gap, Inches(2), term_width, term_height)
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(39, 39, 42)
box2.line.fill.background()

txBox = slide.shapes.add_textbox(term_start + term_width + term_gap + Inches(0.2), Inches(2.2), term_width - Inches(0.4), Inches(2.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Backend Services"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p = tf.add_paragraph()
p.text = "pnpm inventory"
p.font.size = Pt(18)
p.font.name = "Menlo"
p.font.color.rgb = GREEN
p = tf.add_paragraph()
p.text = "localhost:3001"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p = tf.add_paragraph()
p.text = "\nUser + Inventory + Shipping"
p.font.size = Pt(16)
p.font.color.rgb = GRAY

add_subtitle(slide, "json-server with logging shows every request.", top=Inches(5.5), font_size=28, color=YELLOW)

# ============================================================================
# SLIDE 10: The Happy Path
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Happy Path", top=Inches(0.6), font_size=44)
add_bullet_point(slide, "Products show offer badges (Inventory Service)", Inches(1.5), Inches(1.6), WHITE, "1.")
add_bullet_point(slide, "User tier badge shows Pro (User Service)", Inches(1.5), Inches(2.3), WHITE, "2.")
add_bullet_point(slide, "20% discount applied for Pro users", Inches(1.5), Inches(3), WHITE, "3.")
add_bullet_point(slide, "Add to cart (see json-server logs!)", Inches(1.5), Inches(3.7), WHITE, "4.")
add_bullet_point(slide, "Checkout shows shipping options (Shipping Service)", Inches(1.5), Inches(4.4), WHITE, "5.")
add_bullet_point(slide, "Complete purchase", Inches(1.5), Inches(5.1), WHITE, "6.")

add_title(slide, "This works great.", top=Inches(6), font_size=32, color=GREEN)

# ============================================================================
# SLIDE 11: The Testing Problem - Header
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "How Hard Is This To Test?", top=Inches(3), font_size=56)

# ============================================================================
# SLIDE 12: Testing Problem Table - Easy
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Easy Stuff", top=Inches(0.5), font_size=40, color=GREEN)

# Easy scenarios
add_bullet_point(slide, "Happy path", Inches(1), Inches(2), WHITE, "")
add_subtitle(slide, "Just run the app", top=Inches(2), font_size=20, color=GREEN)

add_subtitle(slide, "Only ONE scenario is truly easy.", top=Inches(5.5), font_size=32, color=YELLOW)

# ============================================================================
# SLIDE 13: Testing Problem Table - Annoying
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Annoying Stuff", top=Inches(0.5), font_size=40, color=YELLOW)

add_bullet_point(slide, "Premium user discount", Inches(1), Inches(1.8), WHITE)
add_subtitle(slide, "Edit db.json to change user tier", top=Inches(2.3), font_size=18, color=YELLOW)

add_bullet_point(slide, "Free user pricing", Inches(1), Inches(3), WHITE)
add_subtitle(slide, "Edit db.json again, remember to change back", top=Inches(3.5), font_size=18, color=YELLOW)

add_subtitle(slide, "Possible, but requires manual setup each time.", top=Inches(5.5), font_size=28, color=GRAY)

# ============================================================================
# SLIDE 14: Testing Problem Table - Hard
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Hard Stuff", top=Inches(0.5), font_size=40, color=RED)

add_bullet_point(slide, "Offer ended (0 spots left)", Inches(1), Inches(1.6), WHITE)
add_subtitle(slide, "Edit db.json? Restart server?", top=Inches(2.1), font_size=18, color=RED)

add_bullet_point(slide, "Express shipping unavailable", Inches(1), Inches(2.8), WHITE)
add_subtitle(slide, "Edit db.json to remove express option", top=Inches(3.3), font_size=18, color=RED)

add_bullet_point(slide, "Shipping service error", Inches(1), Inches(4), WHITE)
add_subtitle(slide, "Kill the server mid-test?", top=Inches(4.5), font_size=18, color=RED)

add_bullet_point(slide, "Limited spots (3 left)", Inches(1), Inches(5.2), WHITE)
add_subtitle(slide, "Edit db.json manually every time", top=Inches(5.7), font_size=18, color=RED)

# ============================================================================
# SLIDE 15: Testing Problem - Impossible
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Impossible", top=Inches(0.5), font_size=44, color=RED)

# Big impossible scenario
imp_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(11.333), Inches(2.5))
imp_box.fill.solid()
imp_box.fill.fore_color.rgb = RGBColor(127, 29, 29)  # red-900
imp_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.3), Inches(2), Inches(10.733), Inches(2.1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Offer ends during checkout"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nUser loads page → offer available\nUser clicks checkout → someone takes last spot\nUser clicks pay → ???"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(252, 165, 165)  # red-300
p.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "Edit db.json WHILE the test runs? Time it perfectly?", top=Inches(5), font_size=28, color=GRAY)
add_title(slide, "That's not testing. That's praying.", top=Inches(5.8), font_size=32, color=YELLOW)

# ============================================================================
# SLIDE 16: Parallel Testing Problem
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "50 Tests in Parallel?", top=Inches(1), font_size=48)

add_bullet_point(slide, "All tests hit the same User Service", Inches(2), Inches(2.5), WHITE, "")
add_bullet_point(slide, "All tests hit the same Inventory Service", Inches(2), Inches(3.3), WHITE, "")
add_bullet_point(slide, "All tests hit the same Shipping Service", Inches(2), Inches(4.1), WHITE, "")

add_title(slide, "Shared state = Test conflicts", top=Inches(5.3), font_size=40, color=RED)

# ============================================================================
# SLIDE 17: The Summary Table
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Summary", top=Inches(0.4), font_size=40)

# Create visual summary
add_bullet_point(slide, "Easy:  Happy path only", Inches(1), Inches(1.4), GREEN, "")
add_bullet_point(slide, "Annoying:  Different user tiers", Inches(1), Inches(2.1), YELLOW, "")
add_bullet_point(slide, "Hard:  Service errors, offer states, shipping options", Inches(1), Inches(2.8), RED, "")
add_bullet_point(slide, "Impossible:  Sequences, parallel tests", Inches(1), Inches(3.5), RED, "")

# The more realistic box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(11.333), Inches(1.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(120, 53, 15)  # amber-900
box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.3), Inches(4.7), Inches(10.733), Inches(1.1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "The more realistic your test, the harder it is to set up."
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Edge cases are nearly impossible."
p.font.size = Pt(24)
p.font.color.rgb = YELLOW
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 18: The Solution Tease
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "What if...", top=Inches(2), font_size=48, color=GRAY)
add_title(slide, "Same app. Same code.", top=Inches(3.2), font_size=44)
add_title(slide, "But YOU control the responses.", top=Inches(4.4), font_size=44, color=YELLOW)

# ============================================================================
# SLIDE 19: How Scenarist Works (Preview)
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Scenarist", top=Inches(0.8), font_size=48, color=YELLOW)

# Diagram: json-server still running, but intercepted
add_subtitle(slide, "json-server is still running...", top=Inches(2), font_size=28, color=WHITE)
add_subtitle(slide, "But Scenarist intercepts BEFORE requests reach it.", top=Inches(2.8), font_size=28, color=GREEN)

# Visual
arrow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(4), Inches(9.333), Inches(1.8))
arrow_box.fill.solid()
arrow_box.fill.fore_color.rgb = RGBColor(39, 39, 42)
arrow_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(2.2), Inches(4.2), Inches(8.933), Inches(1.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Next.js → Scenarist → (intercepted) → json-server"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "                                           never reached"
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 20: Next Video
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Next:", top=Inches(2), font_size=36, color=GRAY)
add_title(slide, "One Server, Unlimited Scenarios", top=Inches(3), font_size=48)
add_title(slide, "How Scenarist actually works.", top=Inches(4.5), font_size=36, color=WHITE)

# ============================================================================
# Save the presentation
# ============================================================================
output_path = "video-02-meet-payflow.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
