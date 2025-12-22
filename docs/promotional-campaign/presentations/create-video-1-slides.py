#!/usr/bin/env python3
"""
Generate PowerPoint slides for Video 1: The Testing Gap Nobody Talks About

Run: python3 create-video-1-slides.py
Output: video-01-the-testing-gap.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(24, 24, 27)  # zinc-900
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(239, 68, 68)  # red-500
YELLOW = RGBColor(245, 158, 11)  # amber-500
GREEN = RGBColor(34, 197, 94)  # green-500
BLUE = RGBColor(59, 130, 246)  # blue-500
GRAY = RGBColor(161, 161, 170)  # zinc-400


def add_dark_slide(prs):
    """Add a slide with dark background."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    # Add dark background shape
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide


def add_title(slide, text, top=Inches(2.5), font_size=60, color=WHITE):
    """Add centered title text."""
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_subtitle(slide, text, top=Inches(4), font_size=32, color=GRAY):
    """Add centered subtitle text."""
    txBox = slide.shapes.add_textbox(Inches(1), top, Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_code_block(slide, code, left, top, width, height):
    """Add a code block with monospace font."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(39, 39, 42)  # zinc-800
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Menlo"
    p.font.color.rgb = RGBColor(228, 228, 231)  # zinc-200
    return shape


def add_bullet_point(slide, text, left, top, color=WHITE, icon=""):
    """Add a bullet point with optional icon."""
    txBox = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon}  {text}" if icon else text
    p.font.size = Pt(28)
    p.font.color.rgb = color
    return txBox


# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Testing Gap", top=Inches(2.8), font_size=72)
add_subtitle(slide, "Nobody Talks About", top=Inches(4.2), font_size=40, color=GRAY)

# ============================================================================
# SLIDE 2: The Hook
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Your server tests pass.", top=Inches(1.5), font_size=48)
add_title(slide, "Your frontend tests pass.", top=Inches(2.5), font_size=48)
add_title(slide, "CI is green.", top=Inches(3.5), font_size=48)
add_title(slide, "You ship.", top=Inches(4.5), font_size=48, color=GREEN)
add_title(slide, "Users report a bug.", top=Inches(5.5), font_size=48, color=RED)

# ============================================================================
# SLIDE 3: The Question
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "What went wrong?", top=Inches(3), font_size=64)

# ============================================================================
# SLIDE 4: The Answer
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "You tested both sides separately.", top=Inches(2.5), font_size=44)
add_title(slide, "You never proved they work together.", top=Inches(4), font_size=44, color=YELLOW)

# ============================================================================
# SLIDE 5: Server Test Code
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Server Test", top=Inches(0.5), font_size=36, color=GRAY)

code = """it('returns premium pricing for pro users', async () => {
  vi.mock('./session', () => ({
    getSession: () => ({ userId: '123', tier: 'pro' })
  }));

  const response = await request(app).get('/api/pricing');

  expect(response.body.discount).toBe(20);
});"""

add_code_block(slide, code, Inches(1.5), Inches(1.5), Inches(10), Inches(3.5))
add_title(slide, "Test passes.", top=Inches(5.5), font_size=36, color=GREEN)

# ============================================================================
# SLIDE 6: Server Test Problem
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "What did we prove?", top=Inches(1), font_size=40)
add_bullet_point(slide, "IF session exists correctly...", Inches(2), Inches(2.2), GRAY, "")
add_bullet_point(slide, "AND code reads it correctly...", Inches(2), Inches(3), GRAY, "")
add_bullet_point(slide, "THEN we return the discount.", Inches(2), Inches(3.8), GRAY, "")
add_title(slide, "What we did NOT prove:", top=Inches(5), font_size=36, color=RED)
add_subtitle(slide, "That a real user actually sees it.", top=Inches(5.8), font_size=32, color=WHITE)

# ============================================================================
# SLIDE 7: Frontend Test Code
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Frontend Test", top=Inches(0.5), font_size=36, color=GRAY)

code = """it('displays discount for premium users', () => {
  render(<PricingPage pricing={{ discount: 20 }} />);

  expect(screen.getByText('20% off')).toBeInTheDocument();
});"""

add_code_block(slide, code, Inches(2), Inches(1.8), Inches(9), Inches(2.5))
add_title(slide, "Test passes.", top=Inches(5), font_size=36, color=GREEN)

# ============================================================================
# SLIDE 8: Frontend Test Problem
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "We hardcoded the discount.", top=Inches(2.5), font_size=48)
add_title(slide, "We assumed the server returns it.", top=Inches(3.8), font_size=40, color=GRAY)
add_title(slide, "What if our assumption is wrong?", top=Inches(5.2), font_size=40, color=YELLOW)

# ============================================================================
# SLIDE 9: The Gap - Visual
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "THE GAP", top=Inches(0.5), font_size=48)

# Left box - Server Tests
left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4), Inches(4))
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(127, 29, 29)  # red-900
left_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(3.6), Inches(3.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "SERVER TESTS"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nsupertest / jest / vitest"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\n\nAPI returns JSON"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nNo real browser"
p.font.size = Pt(18)
p.font.color.rgb = RED
p.alignment = PP_ALIGN.CENTER

# Right box - Frontend Tests
right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.833), Inches(1.5), Inches(4), Inches(4))
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(127, 29, 29)  # red-900
right_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(9.033), Inches(1.7), Inches(3.6), Inches(3.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "FRONTEND TESTS"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nRTL / jsdom"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\n\nComponent renders"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nNo real server"
p.font.size = Pt(18)
p.font.color.rgb = RED
p.alignment = PP_ALIGN.CENTER

# Center box - The Gap
center_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(2), Inches(3.333), Inches(3))
center_box.fill.solid()
center_box.fill.fore_color.rgb = RGBColor(120, 53, 15)  # amber-900
center_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(2.933), Inches(2.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "THE GAP"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = YELLOW
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nReal browser\n+ Real server\n+ Controlled APIs"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

add_title(slide, "That's where the bugs hide.", top=Inches(6), font_size=36, color=YELLOW)

# ============================================================================
# SLIDE 10: Why It's Hard
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Why is this so hard?", top=Inches(0.8), font_size=44)
add_title(slide, "Server-side state builds up.", top=Inches(2.2), font_size=36, color=WHITE)

# Flow diagram as text
flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3.2), Inches(11.333), Inches(1.2))
flow_box.fill.solid()
flow_box.fill.fore_color.rgb = RGBColor(39, 39, 42)  # zinc-800
flow_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(10.933), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Login  session  Cart  Checkout  Payment"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "Each step depends on the previous one.", top=Inches(5), font_size=28, color=GRAY)
add_subtitle(slide, "Order matters. State matters.", top=Inches(5.7), font_size=28, color=WHITE)

# ============================================================================
# SLIDE 11: The Impossible Test
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Test: Payment declined for premium user", top=Inches(1), font_size=36)
add_bullet_point(slide, "Real browser making real requests", Inches(1.5), Inches(2.2), WHITE, "1.")
add_bullet_point(slide, "Real server executing your code", Inches(1.5), Inches(3), WHITE, "2.")
add_bullet_point(slide, "Auth0 returning a premium user", Inches(1.5), Inches(3.8), WHITE, "3.")
add_bullet_point(slide, "Stripe returning a declined payment", Inches(1.5), Inches(4.6), WHITE, "4.")

add_title(slide, "You can't ask Stripe to decline for your test.", top=Inches(5.8), font_size=32, color=RED)

# ============================================================================
# SLIDE 12: Testing Pyramid
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "The Testing Pyramid Has a Gap", top=Inches(0.5), font_size=40)

# Draw pyramid layers
# E2E (top)
e2e = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.5), Inches(1.5), Inches(4.333), Inches(1.5))
e2e.fill.solid()
e2e.fill.fore_color.rgb = RGBColor(30, 64, 175)  # blue-800
e2e.line.fill.background()

# The Gap (middle)
gap = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(3.5), Inches(3), Inches(6.333), Inches(1.5))
gap.fill.solid()
gap.fill.fore_color.rgb = RGBColor(120, 53, 15)  # amber-900
gap.line.fill.background()

# Unit (bottom)
unit = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(2.5), Inches(4.5), Inches(8.333), Inches(1.5))
unit.fill.solid()
unit.fill.fore_color.rgb = RGBColor(22, 101, 52)  # green-800
unit.line.fill.background()

# Labels
add_subtitle(slide, "E2E - Can't control APIs", top=Inches(1.6), font_size=20, color=WHITE)
add_subtitle(slide, "THE GAP", top=Inches(3.3), font_size=24, color=YELLOW)
add_subtitle(slide, "Unit - Can't prove integration", top=Inches(4.8), font_size=20, color=WHITE)

# ============================================================================
# SLIDE 13: What You Need
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "What You Need", top=Inches(0.8), font_size=48)

# Three boxes
box_width = Inches(3.5)
box_height = Inches(3)
gap_between = Inches(0.5)
start_left = Inches(1.2)

# Box 1: Real Browser
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_left, Inches(2), box_width, box_height)
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(30, 64, 175)  # blue-800
box1.line.fill.background()

txBox = slide.shapes.add_textbox(start_left + Inches(0.2), Inches(2.3), box_width - Inches(0.4), box_height - Inches(0.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Real Browser"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nPlaywright"
p.font.size = Pt(20)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nNot jsdom"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Box 2: Real Server
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_left + box_width + gap_between, Inches(2), box_width, box_height)
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(22, 101, 52)  # green-800
box2.line.fill.background()

txBox = slide.shapes.add_textbox(start_left + box_width + gap_between + Inches(0.2), Inches(2.3), box_width - Inches(0.4), box_height - Inches(0.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Real Server"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nYour actual app"
p.font.size = Pt(20)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nCode executes"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Box 3: Controlled APIs
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_left + 2 * (box_width + gap_between), Inches(2), box_width, box_height)
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(120, 53, 15)  # amber-900
box3.line.fill.background()

txBox = slide.shapes.add_textbox(start_left + 2 * (box_width + gap_between) + Inches(0.2), Inches(2.3), box_width - Inches(0.4), box_height - Inches(0.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Controlled APIs"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nYou decide responses"
p.font.size = Pt(20)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nAuth0, Stripe, etc."
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

add_title(slide, "Real browser. Real server. Controlled responses.", top=Inches(5.5), font_size=32, color=YELLOW)

# ============================================================================
# SLIDE 14: The Solution Tease (Code)
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "What if you could do this?", top=Inches(0.5), font_size=40, color=GRAY)

code = """test('premium user sees discount', async ({ page, switchScenario }) => {
  await switchScenario('premiumUser');
  await page.goto('/pricing');
  await expect(page.getByText('20% off')).toBeVisible();
});

test('payment declined', async ({ page, switchScenario }) => {
  await switchScenario('paymentDeclined');
  await page.goto('/checkout');
  await page.click('Pay Now');
  await expect(page.getByText('Card declined')).toBeVisible();
});"""

add_code_block(slide, code, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5))

# ============================================================================
# SLIDE 15: Benefits
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "No server restarts.", top=Inches(2), font_size=48, color=GREEN)
add_title(slide, "No config files to edit.", top=Inches(3.2), font_size=48, color=GREEN)
add_title(slide, "No juggling test accounts.", top=Inches(4.4), font_size=48, color=GREEN)
add_title(slide, "Tests run in parallel.", top=Inches(5.6), font_size=48, color=GREEN)

# ============================================================================
# SLIDE 16: Next Video Tease
# ============================================================================
slide = add_dark_slide(prs)
add_title(slide, "Next:", top=Inches(2), font_size=36, color=GRAY)
add_title(slide, "A real payment app.", top=Inches(3), font_size=48)
add_title(slide, "Real Auth0. Real Stripe.", top=Inches(4.2), font_size=40, color=WHITE)
add_title(slide, "Real problems.", top=Inches(5.4), font_size=40, color=YELLOW)

# ============================================================================
# Save the presentation
# ============================================================================
output_path = "video-01-the-testing-gap.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
