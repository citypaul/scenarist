#!/usr/bin/env python3
"""
Generate PowerPoint slides for Video 4: Response Sequences - The Impossible Test

This script creates a presentation showing how Scenarist's response sequences
enable testing scenarios that are impossible with real services.

Usage:
    python3 create-video-4-slides.py

Output:
    video-04-response-sequences.pptx

Requirements:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette (matching Video 1-3)
DARK_BG = RGBColor(24, 24, 27)  # zinc-900
WHITE = RGBColor(255, 255, 255)
ZINC_400 = RGBColor(161, 161, 170)
YELLOW = RGBColor(245, 158, 11)  # amber-500
GREEN = RGBColor(34, 197, 94)  # green-500
RED = RGBColor(239, 68, 68)  # red-500
BLUE = RGBColor(59, 130, 246)  # blue-500
PURPLE = RGBColor(168, 85, 247)  # purple-500


def add_dark_slide(prs):
    """Add a slide with dark background."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    return slide


def add_title(slide, text, top=Inches(0.5), font_size=54, color=WHITE):
    """Add a title text box."""
    left = Inches(0.5)
    width = Inches(12.33)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_subtitle(slide, text, top=Inches(1.5), font_size=28, color=ZINC_400):
    """Add a subtitle text box."""
    left = Inches(0.5)
    width = Inches(12.33)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER


def add_bullet_point(slide, text, left, top, color=WHITE, bullet=""):
    """Add a bullet point."""
    width = Inches(10)
    height = Inches(0.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{bullet} {text}" if bullet else text
    p.font.size = Pt(24)
    p.font.color.rgb = color


def add_code_block(slide, code, left, top, width, height, font_size=14):
    """Add a code block with dark background."""
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(39, 39, 42)  # zinc-800
    shape.line.fill.background()

    # Code text
    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(228, 228, 231)  # zinc-200
    p.font.name = "Menlo"


def add_box(slide, text, left, top, width, height, bg_color, text_color=WHITE, font_size=20):
    """Add a colored box with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    # Add text frame
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_arrow(slide, start_left, start_top, end_left, end_top, color=ZINC_400):
    """Add an arrow connector."""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


def create_presentation():
    """Create the Video 4 presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = add_dark_slide(prs)
    add_title(slide, "Response Sequences", top=Inches(2.5))
    add_subtitle(slide, "The Impossible Test", top=Inches(3.5), font_size=36, color=YELLOW)
    add_subtitle(slide, "Scenarist Video 4", top=Inches(5), font_size=20)

    # Slide 2: The Scenario
    slide = add_dark_slide(prs)
    add_title(slide, "The Scenario", top=Inches(0.5), font_size=44)

    steps = [
        ("1", "User loads product page", "15 items in stock", GREEN),
        ("2", "User adds to cart", "", WHITE),
        ("3", "User fills payment details", "", WHITE),
        ("4", "Meanwhile...", "Someone else buys the last one", RED),
        ("5", "User clicks Pay", "What happens?", YELLOW),
    ]

    top = Inches(1.5)
    for num, action, note, color in steps:
        add_bullet_point(slide, f"{num}. {action}", Inches(1), top, color)
        if note:
            add_bullet_point(slide, note, Inches(5), top, ZINC_400)
        top += Inches(0.7)

    # Slide 3: The Problem
    slide = add_dark_slide(prs)
    add_title(slide, "The Problem", top=Inches(0.5), font_size=44)
    add_subtitle(slide, "Traditional mocks return the same response every time", top=Inches(1.5))

    add_code_block(
        slide,
        'server.use(\n  http.get("/api/inventory", () => {\n    return HttpResponse.json([{ quantity: 15 }]);\n  })\n);',
        Inches(2), Inches(2.5), Inches(9), Inches(2)
    )

    add_subtitle(slide, "Same endpoint. Same response. Every call.", top=Inches(5), color=RED)

    # Slide 4: How Would You Test This?
    slide = add_dark_slide(prs)
    add_title(slide, "How Would You Test This?", top=Inches(0.3), font_size=40)

    options = [
        ("Option 1: Edit db.json mid-test", "Flaky, manual, unreproducible", RED),
        ("Option 2: Two browser sessions", "Impossible to coordinate timing", RED),
        ("Option 3: Mock at function layer", "Not testing real HTTP code", RED),
    ]

    top = Inches(1.5)
    for option, problem, color in options:
        add_bullet_point(slide, option, Inches(0.5), top, WHITE, bullet="❌")
        add_bullet_point(slide, problem, Inches(1), top + Inches(0.5), ZINC_400)
        top += Inches(1.2)

    add_box(slide, "With Scenarist: Define sequence. Run test. Done.",
            Inches(1.5), Inches(5.5), Inches(10), Inches(0.8), GREEN, font_size=24)

    # Slide 5: Introducing Sequences
    slide = add_dark_slide(prs)
    add_title(slide, "Response Sequences", top=Inches(0.5), font_size=44)
    add_subtitle(slide, "Different response on each call", top=Inches(1.3))

    # Visual: Three boxes showing sequence
    add_box(slide, "Call #1\n15 units", Inches(1.5), Inches(2.5), Inches(3), Inches(1.5), GREEN)
    add_box(slide, "Call #2\n0 units", Inches(5.2), Inches(2.5), Inches(3), Inches(1.5), RED)
    add_box(slide, "Call #3+\n0 units", Inches(8.9), Inches(2.5), Inches(3), Inches(1.5), ZINC_400)

    add_subtitle(slide, "Same endpoint → Different responses → Based on call order", top=Inches(5))

    # Slide 6: The Code
    slide = add_dark_slide(prs)
    add_title(slide, "The Sequence Definition", top=Inches(0.3), font_size=40)

    code = '''sequence: {
  responses: [
    // First call (products page): in stock
    { status: 200, body: [{ quantity: 15 }] },

    // Second call (checkout): sold out
    { status: 200, body: [{ quantity: 0 }] },
  ],
  repeat: "last",
}'''

    add_code_block(slide, code, Inches(1.5), Inches(1.3), Inches(10), Inches(4.5), font_size=18)

    add_subtitle(slide, 'repeat: "last" = keep returning final response', top=Inches(6))

    # Slide 7: The Result
    slide = add_dark_slide(prs)
    add_title(slide, "The Result", top=Inches(0.5), font_size=44)

    results = [
        ("✓", "Deterministic", "Same result every time", GREEN),
        ("✓", "Fast", "No timing coordination needed", GREEN),
        ("✓", "Real code paths", "Actual HTTP calls are made", GREEN),
        ("✓", "Zero backend requests", "MSW intercepts everything", GREEN),
    ]

    top = Inches(1.8)
    for icon, title, desc, color in results:
        add_bullet_point(slide, f"{icon} {title}", Inches(1.5), top, color)
        add_bullet_point(slide, desc, Inches(2), top + Inches(0.45), ZINC_400)
        top += Inches(1)

    # Slide 8: Beyond Inventory
    slide = add_dark_slide(prs)
    add_title(slide, "Beyond Inventory", top=Inches(0.5), font_size=44)
    add_subtitle(slide, "Sequences work for anything that changes over time", top=Inches(1.3))

    use_cases = [
        ("Payment Polling", "pending → processing → complete", BLUE),
        ("Retry Logic", "fail → fail → succeed", YELLOW),
        ("Session Expiry", "valid → expired", PURPLE),
        ("State Machines", "any state progression", GREEN),
    ]

    top = Inches(2.3)
    for use_case, progression, color in use_cases:
        add_box(slide, use_case, Inches(1), top, Inches(4), Inches(0.7), color, font_size=18)
        add_bullet_point(slide, progression, Inches(5.5), top + Inches(0.15), WHITE)
        top += Inches(1)

    # Slide 9: Key Takeaway
    slide = add_dark_slide(prs)
    add_title(slide, "Key Takeaway", top=Inches(2))
    add_subtitle(
        slide,
        "Test scenarios that are impossible\nwith real services",
        top=Inches(3.2),
        font_size=36,
        color=YELLOW
    )
    add_subtitle(slide, "Define sequence. Run test. Done.", top=Inches(5), font_size=28)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "video-04-response-sequences.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    create_presentation()
